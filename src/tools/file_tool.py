"""文件解析工具 - 支持多种文档格式解析"""

import os
import re
import tempfile
from typing import Dict, Any, List, Optional, Union
from pathlib import Path

from .base_tool import BaseTool


class FileParserTool(BaseTool):
    """文件解析工具
    
    支持解析多种文档格式：PDF、DOCX、PPTX、XLSX、TXT、CSV等
    """
    
    name = "parse_file"
    description = "Parse multiple user uploaded local files such as PDF, DOCX, PPTX, TXT, CSV, XLSX, DOC, ZIP."
    parameters = {
        "type": "object",
        "properties": {
            "files": {
                "type": "array",
                "items": {"type": "string"},
                "description": "The file name of the user uploaded local files to be parsed."
            }
        },
        "required": ["files"]
    }
    
    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        "pdf", "docx", "doc", "pptx", "ppt",
        "xlsx", "xls", "csv", "txt", "md",
        "json", "xml", "html", "png", "jpg", "jpeg", "webp"
    }
    
    def __init__(
        self, 
        file_root_path: str = "",
        max_content_length: int = 100000,
        cfg: Optional[Dict] = None
    ):
        """初始化文件解析工具
        
        Args:
            file_root_path: 文件根目录
            max_content_length: 最大内容长度
            cfg: 配置字典
        """
        super().__init__(cfg)
        self.file_root_path = file_root_path
        self.max_content_length = max_content_length
    
    async def call(self, params: Union[str, Dict[str, Any]], **kwargs) -> str:
        """异步解析文件
        
        Args:
            params: 包含 files 字段的参数
            
        Returns:
            解析后的文件内容
        """
        params = self._parse_params(params)
        
        try:
            files = params.get("files", params.get("params", {}).get("files", []))
        except:
            return "[parse_file] Invalid request format: Input must be a JSON object containing 'files' field"
        
        if not files:
            return "[parse_file] Invalid request format: 'files' field is required"
        
        # 处理单个文件或多个文件
        if isinstance(files, str):
            files = [files]
        
        results = []
        for file_path in files:
            result = self._parse_single_file(file_path)
            results.append(result)
        
        return "\n\n=======\n\n".join(results)
    
    def _parse_single_file(self, file_path: str) -> str:
        """解析单个文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            解析结果
        """
        # 解析完整路径
        full_path = self._resolve_path(file_path)
        
        if not os.path.exists(full_path):
            return f"[parse_file] File not found: {file_path}"
        
        # 获取文件扩展名
        ext = Path(full_path).suffix.lower().lstrip(".")
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            return f"[parse_file] Unsupported file type: {ext}"
        
        try:
            # 根据扩展名选择解析方法
            if ext == "pdf":
                content = self._parse_pdf(full_path)
            elif ext in ["docx", "doc"]:
                content = self._parse_docx(full_path)
            elif ext in ["pptx", "ppt"]:
                content = self._parse_pptx(full_path)
            elif ext in ["xlsx", "xls"]:
                content = self._parse_excel(full_path)
            elif ext == "csv":
                content = self._parse_csv(full_path)
            elif ext in ["txt", "md", "json", "xml", "html"]:
                content = self._parse_text(full_path)
            elif ext in ["png", "jpg", "jpeg", "webp"]:
                content = self._parse_image(full_path)
            else:
                content = f"Unsupported format: {ext}"
            
            # 截断过长内容
            if len(content) > self.max_content_length:
                content = content[:self.max_content_length] + "\n\n... [Content truncated due to length limit]"
            
            return f"## File: {file_path}\n\n{content}"
            
        except Exception as e:
            return f"[parse_file] Error parsing {file_path}: {str(e)}"
    
    def _resolve_path(self, file_path: str) -> str:
        """解析文件路径
        
        Args:
            file_path: 相对或绝对路径
            
        Returns:
            完整路径
        """
        if os.path.isabs(file_path):
            return file_path
        
        if self.file_root_path:
            return os.path.join(self.file_root_path, file_path)
        
        return file_path
    
    def _parse_pdf(self, file_path: str) -> str:
        """解析 PDF 文件 (增强版：包含表格提取和布局还原)"""
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_content = [f"### Page {i}"]
                    
                    # 1. 尝试提取表格
                    tables = page.extract_tables()
                    if tables:
                        for table in tables:
                            # 过滤掉空表或全空行
                            table = [row for row in table if any(row)]
                            if not table: continue
                            
                            # 转换为 Markdown 格式
                            md_table = self._format_as_md_table(table)
                            page_content.append("\n#### Table Found:\n" + md_table)
                    
                    # 2. 提取文本 (使用 layout=True 尝试保持多栏结构)
                    # 注意：layout=True 需要 pdfplumber 较新版本
                    try:
                        text = page.extract_text(layout=True)
                    except:
                        text = page.extract_text()
                        
                    if text:
                        page_content.append("\n#### Text Content:\n" + text)
                    
                    text_parts.append("\n".join(page_content))
            
            return "\n\n---\n\n".join(text_parts) if text_parts else "No content found in PDF."
            
        except ImportError:
            return "[PDF parsing requires pdfplumber: pip install pdfplumber]"
        except Exception as e:
            return f"PDF parsing error: {str(e)}"

    def _format_as_md_table(self, table: List[List[Optional[str]]]) -> str:
        """将提取的表格列表转换为 Markdown 表格字符"""
        if not table: return ""
        
        # 处理 None 值
        clean_table = [[(str(cell) if cell is not None else "") for cell in row] for row in table]
        
        if not clean_table: return ""
        
        headers = clean_table[0]
        rows = clean_table[1:]
        
        md = "| " + " | ".join(headers) + " |\n"
        md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        
        for row in rows:
            # 确保每行长度对齐
            if len(row) < len(headers):
                row.extend([""] * (len(headers) - len(row)))
            md += "| " + " | ".join(row[:len(headers)]) + " |\n"
            
        return md
    
    def _parse_docx(self, file_path: str) -> str:
        """解析 Word 文档"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            return "\n\n".join(paragraphs) if paragraphs else "No text content found in document."
            
        except ImportError:
            return "[DOCX parsing requires python-docx: pip install python-docx]"
        except Exception as e:
            return f"DOCX parsing error: {str(e)}"
    
    def _parse_pptx(self, file_path: str) -> str:
        """解析 PowerPoint 文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if slide_content:
                    slides_text.append(f"### Slide {i}\n" + "\n".join(slide_content))
            
            return "\n\n".join(slides_text) if slides_text else "No text content found in presentation."
            
        except ImportError:
            return "[PPTX parsing requires python-pptx: pip install python-pptx]"
        except Exception as e:
            return f"PPTX parsing error: {str(e)}"
    
    def _parse_excel(self, file_path: str) -> str:
        """解析 Excel 文件"""
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            sheets_text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 转换为 Markdown 表格
                md_table = df.head(100).to_markdown(index=False)  # 限制行数
                sheets_text.append(f"### Sheet: {sheet_name}\n\n{md_table}")
            
            return "\n\n".join(sheets_text)
            
        except ImportError:
            return "[Excel parsing requires pandas and openpyxl: pip install pandas openpyxl]"
        except Exception as e:
            return f"Excel parsing error: {str(e)}"
    
    def _parse_csv(self, file_path: str) -> str:
        """解析 CSV 文件"""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            
            # 转换为 Markdown 表格（限制行数）
            return df.head(100).to_markdown(index=False)
            
        except ImportError:
            return "[CSV parsing requires pandas: pip install pandas]"
        except Exception as e:
            return f"CSV parsing error: {str(e)}"
    
    def _parse_text(self, file_path: str) -> str:
        """解析纯文本文件"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            with open(file_path, "r", encoding="gbk") as f:
                return f.read()
        except Exception as e:
            return f"Text parsing error: {str(e)}"

    def _parse_image(self, file_path: str) -> str:
        """解析图像文件"""
        try:
            from .multimodal_tool import ImageAnalysisTool
            vlm_tool = ImageAnalysisTool()
            return vlm_tool.call({"image_path": [file_path], "prompt": "Identify and describe all details in this image, including text and objects."})
        except Exception as e:
            return f"Image parsing error: {str(e)}"

